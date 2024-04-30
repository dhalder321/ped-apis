import json, os, time
from datetime import datetime
import base64
import logging
from pathlib import Path
from common.db import DBManager
from datetime import datetime, timezone
from docx import Document
from pptx import Presentation
from htmldocx import HtmlToDocx
from common.s3File import uploadFile 
from common.gmail import sendCompanyEmail 
from random import randint

class Utility:
  

  ##################################################
  #Common terms in the prompts to replace 
  ###

  # System role: {{SYSTEM_ROLE}}
  # Topic: {{TOPIC}}
  # Summary: {{SUMMARY}}
  # Outline: {{OUTLINE}}
  
  ##################################################

  TOPIC2SUMMARY_PROMPT_TYPE = 'TOPIC2SUMMARY'
  TEXT2TOPICOUTLINE_PROMPT_TYPE = 'TEXT2TOPICOUTLINE'
  TEXTOFTOPICOUTLINE_PROMPT_TYPE = 'TEXTOFTOPICOUTLINE'
  QUICK_TEXT_PROMPT_TYPE = 'QUICK_TEXT'

  ESSAY_MODEL_OVERRIDE_PROMPT_TYPE = "ESSAY_MODEL_OVERRIDE_PROMPT"                         
  ESSAY_MODEL_SUBHEADING_PROMPT_TYPE = "ESSAY_MODEL_SUBHEADING_PROMPT"                        
  MEDIUM_ESSAY_MODEL_OVERRIDE_PROMPT_TYPE =  "MEDIUM_ESSAY_MODEL_OVERRIDE_PROMPT"                 
  MEDIUM_ESSAY_MODEL_SUBHEADING_PROMPT_TYPE =  "MEDIUM_ESSAY_MODEL_SUBHEADING_PROMPT"               
  SHORT_ESSAY_MODEL_OVERRIDE_PROMPT_TYPE =  "SHORT_ESSAY_MODEL_OVERRIDE_PROMPT"                  
  SHORT_ESSAY_MODEL_SUBHEADING_PROMPT_TYPE =  "SHORT_ESSAY_MODEL_SUBHEADING_PROMPT"               

  TRANSFORM_ANNOTATED_BIBLIOGRAPHY_PROMPT_TYPE =  "TRANSFORM_ANNOTATED_BIBLIOGRAPHY"         
  TRANSFORM_CRITICAL_RESPONSE_ESSAY_PROMPT_TYPE =  "TRANSFORM_CRITICAL_RESPONSE_ESSAY"        
  TRANSFORM_CRITICAL_RESPONSE_ESSAY_INSTRUCTION_PROMPT_TYPE =   "TRANSFORM_CRITICAL_RESPONSE_ESSAY_INSTRUCTION"                                        
  TRANSFORM_ESSAY_EXPANSION_ASSIGNMENT_PROMPT_TYPE =    "TRANSFORM_ESSAY_EXPANSION_ASSIGNMENT"   
  TRANSFORM_REFLECTION_PAPER_PROMPT_TYPE =    "TRANSFORM_REFLECTION_PAPER"             
  TRANSFORM_REFLECTION_PAPER_INSTRUCTION_PROMPT_TYPE =  "TRANSFORM_REFLECTION_PAPER_INSTRUCTION"   
  TRANSFORM_RESEARCH_PROPOSALS_PROMPT_TYPE =   "TRANSFORM_RESEARCH_PROPOSALS"            
  TRANSFORM_SEMINAR_DISCUSSION_POINTS_PROMPT_TYPE =  "TRANSFORM_SEMINAR_DISCUSSION_POINTS"      
  TRANSFORM_STUDY_GUIDE_PROMPT_TYPE =   "TRANSFORM_STUDY_GUIDE"    

  TRANSFORM_PPT_CONTENT_LISTANDHEADINGONLY_INSTRUCTION_PROMPT_TYPE = "TRANSFORM_PPT_CONTENT_LISTANDHEADINGONLY_INSTRUCTION"
  TRANSFORM_PPT_CONTENT_LISTANDTEXTONLY_INSTRUCTION_PROMPT_TYPE = "TRANSFORM_PPT_CONTENT_LISTANDTEXTONLY_INSTRUCTION"
  TRANSFORM_PPT_CONTENT_TEXTONLY_INSTRUCTION_PROMPT_TYPE = "TRANSFORM_PPT_CONTENT_TEXTONLY_INSTRUCTION"
  TRANSFORM_PPT_FULLTEXT_PROMPT_TYPE = "TRANSFORM_PPT_FULLTEXT"
  TRANSFORM_PPT_JSON_LISTANDHEADING_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_LISTANDHEADING_FORMAT"
  TRANSFORM_PPT_JSON_LISTANDHEADING_NOTES_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_LISTANDHEADING_NOTES_FORMAT"
  TRANSFORM_PPT_JSON_LISTANDTEXT_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_LISTANDTEXT_FORMAT"
  TRANSFORM_PPT_JSON_LISTANDTEXT_NOTES_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_LISTANDTEXT_NOTES_FORMAT"
  TRANSFORM_PPT_JSON_TEXT_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_TEXT_FORMAT"
  TRANSFORM_PPT_JSON_TEXT_NOTES_FORMAT_PROMPT_TYPE = "TRANSFORM_PPT_JSON_TEXT_NOTES_FORMAT"
  TRANSFORM_PPT_NOTES_INSTRUCTION_PROMPT_TYPE = "TRANSFORM_PPT_NOTES_INSTRUCTION"
  TRANSFORM_PPT_SUMMARY_PROMPT_TYPE = "TRANSFORM_PPT_SUMMARY"
  TRANSFORM_PPT_OVERRIDE_PROMPT_TYPE = "TRANSFORM_PPT_OVERRIDE_PROMPT"
  TRANSFORM_PPT_SUBHEADING_PROMPT_TYPE = "TRANSFORM_PPT_SUBHEADING_PROMPT"

  TRANSFORM_QUIZ_PROMPT_TYPE = "TRANSFORM_QUIZ"
  QUIZ_2_QUESTION_DOCUMENT_PROMPT_TYPE = "QUIZ_2_QUESTION_DOCUMENT_PROMPT"
  QUIZ_2_QnA_DOCUMENT_PROMPT_TYPE = "QUIZ_2_QnA_DOCUMENT_PROMPT"

  SLIDE_NOTE_FOR_VIDEO_PROMPT_TYPE = "SLIDE_NOTE_FOR_VIDEO" 

  PROMPT_TYPE2FILE_NAME = {
    TOPIC2SUMMARY_PROMPT_TYPE: "TOPIC2SUMMARY.txt",
    TEXT2TOPICOUTLINE_PROMPT_TYPE: "TEXT2TOPICOUTLINE.txt",
    TEXTOFTOPICOUTLINE_PROMPT_TYPE: "TEXTOFTOPICOUTLINE.txt",
    QUICK_TEXT_PROMPT_TYPE: "QUICK_TEXT.txt",

    ESSAY_MODEL_OVERRIDE_PROMPT_TYPE  :  "ESSAY_MODEL_OVERRIDE_PROMPT.txt",                         
    ESSAY_MODEL_SUBHEADING_PROMPT_TYPE  :  "ESSAY_MODEL_SUBHEADING_PROMPT.txt",                        
    MEDIUM_ESSAY_MODEL_OVERRIDE_PROMPT_TYPE  :   "MEDIUM_ESSAY_MODEL_OVERRIDE_PROMPT.txt",                 
    MEDIUM_ESSAY_MODEL_SUBHEADING_PROMPT_TYPE  :   "MEDIUM_ESSAY_MODEL_SUBHEADING_PROMPT.txt",               
    SHORT_ESSAY_MODEL_OVERRIDE_PROMPT_TYPE  :   "SHORT_ESSAY_MODEL_OVERRIDE_PROMPT.txt",                  
    SHORT_ESSAY_MODEL_SUBHEADING_PROMPT_TYPE  :   "SHORT_ESSAY_MODEL_SUBHEADING_PROMPT.txt",  
    
    TRANSFORM_ANNOTATED_BIBLIOGRAPHY_PROMPT_TYPE  :   "TRANSFORM_ANNOTATED_BIBLIOGRAPHY.txt",        
    TRANSFORM_CRITICAL_RESPONSE_ESSAY_PROMPT_TYPE  :   "TRANSFORM_CRITICAL_RESPONSE_ESSAY.txt",       
    TRANSFORM_CRITICAL_RESPONSE_ESSAY_INSTRUCTION_PROMPT_TYPE  :    "TRANSFORM_CRITICAL_RESPONSE_ESSAY_INSTRUCTION.txt",                                       
    TRANSFORM_ESSAY_EXPANSION_ASSIGNMENT_PROMPT_TYPE  :     "TRANSFORM_ESSAY_EXPANSION_ASSIGNMENT.txt",  
    TRANSFORM_REFLECTION_PAPER_PROMPT_TYPE  :     "TRANSFORM_REFLECTION_PAPER.txt",            
    TRANSFORM_REFLECTION_PAPER_INSTRUCTION_PROMPT_TYPE  :   "TRANSFORM_REFLECTION_PAPER_INSTRUCTION.txt",  
    TRANSFORM_RESEARCH_PROPOSALS_PROMPT_TYPE  :    "TRANSFORM_RESEARCH_PROPOSALS.txt",           
    TRANSFORM_SEMINAR_DISCUSSION_POINTS_PROMPT_TYPE  :   "TRANSFORM_SEMINAR_DISCUSSION_POINTS.txt",     
    TRANSFORM_STUDY_GUIDE_PROMPT_TYPE  :    "TRANSFORM_STUDY_GUIDE.txt",

    TRANSFORM_PPT_CONTENT_LISTANDHEADINGONLY_INSTRUCTION_PROMPT_TYPE : "TRANSFORM_PPT_CONTENT_LISTANDHEADINGONLY_INSTRUCTION.txt",
    TRANSFORM_PPT_CONTENT_LISTANDTEXTONLY_INSTRUCTION_PROMPT_TYPE : "TRANSFORM_PPT_CONTENT_LISTANDTEXTONLY_INSTRUCTION.txt",
    TRANSFORM_PPT_CONTENT_TEXTONLY_INSTRUCTION_PROMPT_TYPE : "TRANSFORM_PPT_CONTENT_TEXTONLY_INSTRUCTION.txt",
    TRANSFORM_PPT_FULLTEXT_PROMPT_TYPE : "TRANSFORM_PPT_FULLTEXT.txt",
    TRANSFORM_PPT_JSON_LISTANDHEADING_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_LISTANDHEADING_FORMAT.txt",
    TRANSFORM_PPT_JSON_LISTANDHEADING_NOTES_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_LISTANDHEADING_NOTES_FORMAT.txt",
    TRANSFORM_PPT_JSON_LISTANDTEXT_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_LISTANDTEXT_FORMAT.txt",
    TRANSFORM_PPT_JSON_LISTANDTEXT_NOTES_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_LISTANDTEXT_NOTES_FORMAT.txt",
    TRANSFORM_PPT_JSON_TEXT_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_TEXT_FORMAT.txt",
    TRANSFORM_PPT_JSON_TEXT_NOTES_FORMAT_PROMPT_TYPE : "TRANSFORM_PPT_JSON_TEXT_NOTES_FORMAT.txt",
    TRANSFORM_PPT_NOTES_INSTRUCTION_PROMPT_TYPE : "TRANSFORM_PPT_NOTES_INSTRUCTION.txt",
    TRANSFORM_PPT_SUMMARY_PROMPT_TYPE : "TRANSFORM_PPT_SUMMARY.txt",
    TRANSFORM_PPT_OVERRIDE_PROMPT_TYPE : "TRANSFORM_PPT_OVERRIDE_PROMPT.txt",
    TRANSFORM_PPT_SUBHEADING_PROMPT_TYPE : "TRANSFORM_PPT_SUBHEADING_PROMPT.txt",

    TRANSFORM_QUIZ_PROMPT_TYPE : "TRANSFORM_QUIZ.txt",
    QUIZ_2_QUESTION_DOCUMENT_PROMPT_TYPE : "QUIZ_2_QUESTION_DOCUMENT_PROMPT.txt",
    QUIZ_2_QnA_DOCUMENT_PROMPT_TYPE : "QUIZ_2_QnA_DOCUMENT_PROMPT.txt",

    SLIDE_NOTE_FOR_VIDEO_PROMPT_TYPE : "SLIDE_NOTE_FOR_VIDEO.txt"               
  }

  # Local_Location = 'C:\openai-sdk\ped-apis'
  Local_Location = './'
  Efs_Path = '/mnt/ped'
  EFS_LOCATION = Local_Location

  PED_EMAIL_SENDER = 'contact@pioneereducationtech.com'

  DOC2SFDT_LAMBDA_FUNCTION_NAME = 'arn:aws:lambda:us-east-2:464311745778:function:ped-doc2sfdt'

  S3BUCKE_NAME = 'pedbuc'
  S3OBJECT_NAME_FOR_USER_FILES = 'user-files'
  S3OBJECT_NAME_FOR_PROMPT_FILES = 'prompts'
  S3OBJECT_NAME_FOR_TEMPORARY_FILES = 'temp'

  TRASFORMATION_USER_ROLE = "You are a seasoned and experienced academician"
  PROMPT_EXTENSION_4_HTML_OUTPUT = "Generate the output strictly and strictly in HTML format with at minimum doctype, html, head and body tags and other basic HTML tags. Do NOT add any meta tags."

  VIDEO_GENERATION_BACKGROUND_MUSIC_FILE_NAME = "learning-video-background-music.mp3"
  PPT2IMAGE_GENERATION_LAMBDA_FUNCTION_NAME = "ped-getimagesfromppt"
  BASIC_PPT_TEXT_TEMPLATE_FILE_NAME = "BasicPresentationTextTemplate.pptx"
  BASIC_PPT_LIST_TEMPLATE_FILE_NAME = "BasicPresentationListTemplate.pptx"

  ###################################################
  #              Environment variables
  ####################################################
  ENVIRONMENT = "dev"
  PROMPT_LOCATION = "local"  # or s3

  # CORS allowed origin
  CORS_ALLOWED_ORIGIN = "http://localhost:3000" 

  # Windows PPT to Image generation API URL
  PPT_2_IMAGE_GENERATION_API_URL = "ped-getimagesfromppt"

  ###################################################


  ###################################################
  #              TABLE NAMES
  ####################################################
  #Environment: dev
  #Table names:
  USER_TABLE_NAME="ped-users"
  USERFILES_TABLE_NAME = "ped-userfiles" 
  USERACTIVITY_TABLE_NAME = "ped-useractivity" 
  TEMPUSER_TABLE_NAME = "ped-temporaryusers" 

  #Environment: test
  #Table names:
  TEST_USER_TABLE_NAME="ped-users-test"
  TEST_USERFILES_TABLE_NAME = "ped-userfiles-test" 
  TEST_USERACTIVITY_TABLE_NAME = "ped-useractivity-test"
  TEST_TEMPUSER_TABLE_NAME = "ped-temporaryusers-test" 


  #Environment: prod
  #Table names:
  PROD_USER_TABLE_NAME="ped-users-prod"
  PROD_USERFILES_TABLE_NAME = "ped-userfiles-prod" 
  PROD_USERACTIVITY_TABLE_NAME = "ped-useractivity-prod"
  PROD_TEMPUSER_TABLE_NAME = "ped-temporaryusers-prod" 
  ######################################################

  @staticmethod
  def initiate():
    #if Utility.PROMPT_LOCATION != 'dev':
    Utility.EFS_LOCATION = Utility.Efs_Path


  @staticmethod
  def generateResponse(responseCode, bodyJson, origin=None, headers=None):
      
      try:
      
        if headers is None:
            headers= {}

        headers['Content-Type'] = 'application/json'
        #add CORS headers
        headers['Access-Control-Allow-Headers'] = 'Content-Type'
        headers['Access-Control-Allow-Origin'] = Utility.CORS_ALLOWED_ORIGIN if origin is None else origin
        headers['Access-Control-Allow-Methods'] = 'OPTIONS,POST,GET'

        return {
                    'statusCode': responseCode,
                    'body': json.dumps(bodyJson),
                    'headers': headers
                }
      except Exception as e:
        logging.error(e)
        return None
  
  @staticmethod
  def saveBase64FileInLocal(filePath, fileContentBase64):
    try:

      # check if folder exists, create it if not present
      path = Path(filePath)
      fileLocation = str(path.parent)
      isExist = os.path.exists(fileLocation)
      if not isExist:    
          os.makedirs(fileLocation)  

      #write the file
      with open(filePath,"wb") as f:
            f.write(base64.b64decode(fileContentBase64.encode("utf-8")))
      return True
    
    except Exception as e:
      logging.error(e)
      return False

  @staticmethod
  def uploadDocumentinHTMLtoS3(htmlContent, fileName, localFileLocation, s3FfilePath):

    #check for valid filename and file location
    if fileName is None or localFileLocation is None:
      return None
       
    #check for file location existance
    isExist = os.path.exists(localFileLocation)
    if not isExist:    
      os.makedirs(localFileLocation)

    localFilePath = str(Path(localFileLocation, fileName))

    document = Document()
    parser =  HtmlToDocx()
    parser.add_html_to_document(htmlContent, document)
    document.save(localFilePath)

    #upload to S3
    s3fileLocation = Utility.S3OBJECT_NAME_FOR_USER_FILES + "/" + Utility.ENVIRONMENT
    s3filePath = s3fileLocation + s3FfilePath
    print("s3filePath:" + s3filePath)
    presignedURL = uploadFile(localFilePath, Utility.S3BUCKE_NAME, s3filePath)
    
    return presignedURL
  
  @staticmethod
  def uploadDocumentinBase64toS3(fileContent, fileName, localFileLocation, s3FfilePath):
    
    #check for valid filename and file location
    if fileName is None or localFileLocation is None:
      return None
       
    #check for file location existance
    isExist = os.path.exists(localFileLocation)
    if not isExist:    
      os.makedirs(localFileLocation)

    localFilePath = str(Path(localFileLocation, fileName))

    if Utility.saveBase64FileInLocal(localFilePath, fileContent) == False:
      return None
    
    #upload to S3
    s3fileLocation = Utility.S3OBJECT_NAME_FOR_USER_FILES + "/" + Utility.ENVIRONMENT
    s3filePath = s3fileLocation + s3FfilePath
    print("s3filePath:" + s3filePath)
    presignedURL = uploadFile(localFilePath, Utility.S3BUCKE_NAME, s3filePath)
    
    return presignedURL
  

  @staticmethod
  def uploadDocumentinTexttoS3(textContent, fileName, localFileLocation, s3FfilePath):

    #check for valid filename and file location
    if fileName is None or localFileLocation is None:
      return None
       
    #check for file location existance
    isExist = os.path.exists(localFileLocation)
    if not isExist:    
      os.makedirs(localFileLocation)

    localFilePath = str(Path(localFileLocation, fileName))

    document = Document()
    document.add_paragraph(textContent)
    document.save(localFilePath)

    #upload to S3
    s3fileLocation = Utility.S3OBJECT_NAME_FOR_USER_FILES + "/" + Utility.ENVIRONMENT
    s3filePath = s3fileLocation + s3FfilePath
    print("s3filePath:" + s3filePath)
    presignedURL = uploadFile(localFilePath, Utility.S3BUCKE_NAME, s3filePath)
    
    return presignedURL
  
  @staticmethod
  def uploadPPTinJSONtoS3(slidejson, fileName, localFileLocation, s3FfilePath, layoutType):

    #check for file location existance
    isExist = os.path.exists(localFileLocation)
    if not isExist:    
      os.makedirs(localFileLocation)
    
    r = json.loads(slidejson)

#         {
#     "title": "Understanding Globalization", 
#     "slides": [
#         {
#             "heading": "Concept and Impact of Globalization", 
#             "content": [
#                 "Globalization is a broad process fostering global interdependence. It involves the diffusion of knowledge, ideas, norms, and human interactions globally.",
#                 "It has revolutionized communication, trade, technology, and even culture. Globalization has the potential to make societies richer through trade and enables knowledge and ideas to be shared.",
#                 "However, globalization can create challenges such as income inequality. Not all individuals, industries, and regions are equally well-equipped to seize the benefits.", 
#                 "It can also lead to cultural assimilation potentially eroding the cultural diversity and identities of affected societies, leading to a form of cultural homogenization."
#             ]
#         }, 
#         {
#             "heading": "Globalization in The Modern World",
#             "content": [
#                 "Today, globalization is primarily economic involving the transnational circulation of goods and capital. It has brought unparalleled economic growth and raised worldwide standards of living.",
#                 "Modern Information and Communication Technologies (ICT) are accelerating the globalization process, making the world more interconnected.",
#                 "Globalization is also seen in the political sphere with the formation of international organizations like the UN. Such organizations aim to deal with global challenges collectively.",
#                 "However, recent years have seen new challenges to globalization. Anti-globalization movements and populist sentiments are rising due to concerns over national sovereignty and economic distribution."
#             ]
#         }
#     ]
# }

    title = r["title"]
    slide_data = r["slides"]


    # """ Ref for slide types:  
    # 0 ->  title and subtitle 
    # 1 ->  title and content 
    # 2 ->  section header 
    # 3 ->  two content 
    # 4 ->  Comparison 
    # 5 ->  Title only  
    # 6 ->  Blank 
    # 7 ->  Content with caption 
    # 8 ->  Pic with caption 
    # """
    
    if layoutType == 'TEXT':
      basePPTFilePath = str(Path(Utility.EFS_LOCATION, Utility.BASIC_PPT_TEXT_TEMPLATE_FILE_NAME))
      prs = Presentation(basePPTFilePath)
      print("Text template picked up.")
    else:
      basePPTFilePath = str(Path(Utility.EFS_LOCATION, Utility.BASIC_PPT_LIST_TEMPLATE_FILE_NAME))
      prs = Presentation(basePPTFilePath)
      print("List template picked up.")
      
    if title is not None:
        title_slide = prs.slides[0] #.add_slide(prs.slide_layouts[0])
        title_text = title_slide.shapes.title
        title_text.text = title
        title_slide.placeholders[1].text =  "Pioneer Education Tech"

    slideNo = 1
    for slide in slide_data: #iteration on slides

      new_slide = prs.slides[slideNo]  #.add_slide(slide_layout)
  
      if 'heading' in slide.keys() and slide['heading']:
          ttl = new_slide.shapes.title
          ttl.text = slide['heading']

      if layoutType == 'TEXT':
        if 'content' in slide.keys() and slide['content']:
          content = slide['content']
      elif layoutType == 'LIST':
        if 'content' in slide.keys() and slide['content']:
          content = ''
          for c in slide['content']:
            content += c + '\n\n'

      if 'notes' in slide.keys() and slide['notes']:
          notes_slide = new_slide.notes_slide
          text_frame = notes_slide.notes_text_frame
          text_frame.text = slide['notes']

      if content is not None and content != "":
          shapes = new_slide.shapes
          body_shape = shapes.placeholders[1]
          tf = body_shape.text_frame
          tf.text = content
          # tf.fit_text(font_family="Arial", max_size=18)

      slideNo += 1

    #delete rest of the slides
    slideList = prs.slides._sldIdLst 
    slides = list(slideList)
    for i in range(len(slides) - 1,slideNo,-1):
        slide = prs.slides[i-1]
        slideList.remove(slides[i-1])

    #convert the pptx into a byte stream
    # deck_name = title + ".pptx"
    localFilePath = str(Path(localFileLocation, fileName)) 
    prs.save(localFilePath)
    
    # with io.BytesIO() as out:
    #     prs.save(out)
    #     out.seek(0)
    #     s3.upload_fileobj(out, s3bucket, deck_name)

    #upload to S3
    s3fileLocation = Utility.S3OBJECT_NAME_FOR_USER_FILES + "/" + Utility.ENVIRONMENT
    s3filePath = s3fileLocation + s3FfilePath
    print("s3filePath:" + s3filePath)
    presignedURL = uploadFile(localFilePath, Utility.S3BUCKE_NAME, s3filePath)
    
    return presignedURL
  
  @staticmethod
  def uploadQuizinJSONtoS3(quizJson, fileName, localFileLocation, s3FfilePath):

    if localFileLocation is None or fileName is None:
      return None

    isExist = os.path.exists(localFileLocation)
    if not isExist:    
      os.makedirs(localFileLocation)
    
    localFilePath = str(Path(localFileLocation, fileName)) 
    with open(localFilePath, 'w') as file:
      file.write(quizJson)

    #upload to S3
    s3fileLocation = Utility.S3OBJECT_NAME_FOR_USER_FILES + "/" + Utility.ENVIRONMENT
    s3filePath = s3fileLocation + s3FfilePath
    print("s3filePath:" + s3filePath)
    presignedURL = uploadFile(localFilePath, Utility.S3BUCKE_NAME, s3filePath)
    
    return presignedURL
        
  @staticmethod
  def curtailObject4Logging(body, field):
     
    bodyCurtailed = body.copy()
    if field in bodyCurtailed:
      # bodyCurtailed.pop(field, None)
      bodyCurtailed[field] = "..."

    return bodyCurtailed 
      
  @staticmethod
  def logUserActivity(body, methodName):

    user_id = body['userid'] if 'userid' in body else None
    tran_id_ = body['transactionId'] if 'transactionId' in body else None
    requestTime = body['requesttimeinUTC'] if 'requesttimeinUTC' in body else ""
    
    if user_id is None and methodName not in ("signupNewUser", "loginUserWithemail", "loginUserWithAccessKey", "generateEmailVerificationCode", "verifyEmailVerificationCode", "getQuizJSON"):
      raise ValueError("userid not sent in request")
    elif user_id is None and methodName in ("signupNewUser", "loginUserWithemail", "loginUserWithAccessKey", "generateEmailVerificationCode", "verifyEmailVerificationCode", "getQuizJSON"):
      user_id = "-1"
    else:
      if not user_id.isdigit():
        raise ValueError("User id sent in request is not a valid integer")

    if tran_id_ is None:
      raise ValueError("transactionid not sent in request")   
    
    return DBManager.addRecordInDynamoTableWithAutoIncrKey(DBTables.UserActivity_Table_Name, \
                                                    "staticIndexColumn", "activityid", \
                                                    "staticIndexColumn-activityid-index", \
                                                     {
                                                        "userid": int(user_id),
                                                        "transactionid": tran_id_,
                                                        "staticIndexColumn": 99,
                                                        "requestTimeInUTC": requestTime,
                                                        "apiMethodName": methodName,
                                                        "activityType": "APIInvoke",
                                                        "requestBody": json.dumps(body)
                                                      })
  
  @staticmethod
  def updateUserActivity(activityId, userid, response):
     if activityId is None or userid is None:
        return None
     
     return DBManager.updateRecordInDynamoTable(DBTables.UserActivity_Table_Name, \
                                              "activityid", activityId, \
                                              "userid", userid, \
                                                     {
                                                        "responseTimeInUTC": datetime.now().replace(tzinfo=timezone.utc).strftime("%m/%d/%Y %H:%M:%S"),
                                                        "responseStatus": "success" if response['statusCode'] == 200 else "failure", 
                                                        "responseBody": json.dumps(response)
                                                      })
  
  @staticmethod
  def handlePriorTransactionIds(userId, tran_IDs):

    # look for existing transaction id with completion status
    # if found, return the response
    if tran_IDs is None or len(tran_IDs) <= 0:
      return
    
    trans = tran_IDs.split(",")

    for t in trans:
      if t == "":
        continue

      records = DBManager.getDBItemByIndex(DBTables.UserActivity_Table_Name, \
                                           "transactionid", "transactionid-index", \
                                            t)
      
      print(str(records))
      # if record present, then return response or wait for it to complete
      # continue with next tran id if 500 error occured
      if records is not None and len(records) == 1:

        # match with user id - essential for data privacy issues
        if 'userid' in records[0] and str(records[0]['userid']) != str(userId):
          continue

        if "responseStatus" in records[0] and records[0]["responseStatus"] == "success":
          if "responseBody" in records[0] and records[0]["responseBody"] != "":
                return json.loads(records[0]["responseBody"])
          else: # no response found
            continue
        # API error 
        elif "responseStatus" in records[0] and records[0]["responseStatus"] != "success":
          continue
        
        # start a loop for 3 times to get the running process to complete 
        if "responseStatus" not in records[0]:

          for n in range(3):
            if "responseStatus" in records[0] and records[0]["responseStatus"] == "success":

              # return the response
              if "responseBody" in records[0] and records[0]["responseBody"] != "":
                return json.loads(records[0]["responseBody"])

              else: # success but no response found, db updated failed somehow
                continue
              
            # record found but not response status or API failed
            elif "responseStatus" in records[0] and records[0]["responseStatus"] != "success": 
              continue
            
            # wait for 15 sec
            time.sleep(15)

            records = DBManager.getDBItemByIndex(DBTables.UserActivity_Table_Name, \
                                            "transactionid", "transactionid-index", \
                                              t)
      else: # no record or more than one record found
        continue


    logging.debug(msg="*******No successful transaction found for transaction ids::" + str(tran_IDs))      

    return None

  @staticmethod
  def formatLogMessage(tran_id, userId, message):

    msg = "Transaction ID:\"{transactionId}\"- User ID: \"{userID}\"- Message: \"{message}\""

    if tran_id is None:
      tran_id = ''
    if userId is None:
      userid = ''
    if message is None:
      message = ''

    return msg.format(transactionId=tran_id, userID = userId, message=message)

  @staticmethod
  def randomNumberOfNDigits(n):
    range_start = 10**(n-1)
    range_end = (10**n)-1
    return randint(range_start, range_end)
  
  @staticmethod
  def sendEmailConfirmationCodeInEmail(receiverEmail):

    # generate a 6 digit code that expires in 10 mins
    randomNumber = Utility.randomNumberOfNDigits(6)

    subject = str(randomNumber) + " Pioneer Education - your email verification code"

    body = '''
      Hello there!

            {{CODE}} is your email verification code.
        
        Please use above code to verify your email address during signup process. With this verification, you will confirm your email address with Pioneer Education Tech corporation. 

        If you dont recognize this email, please drop us an email - contact@pioneereducationtech.com 

    '''
    body = body.replace('{{CODE}}', str(randomNumber))

    if sendCompanyEmail(Utility.PED_EMAIL_SENDER, receiverEmail, subject, body) is None:
      return False
    
    # log the email verification code
    retVal = DBManager.addRecordInDynamoTableWithAutoIncrKey(DBTables.TempUser_Table_Name, "staticIndexColumn",
                                                    "tempUserId", "staticIndexColumn-tempUserId-index",
                                                    {
                                                      "staticIndexColumn": 99,
                                                      "email": receiverEmail,
                                                      "emailVerificationCode": str(randomNumber),
                                                      "verificationCodeGenerationTime" : datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
                                                      "verified": "n",
                                                    })
    if retVal is not None:
      return True
    else:
      return False


  @staticmethod
  def verifyEmailVerificationCode(receiverEmail, emailCode):

    # get the code from DB
    record = DBManager.getDBItemByIndex(DBTables.TempUser_Table_Name, "email",
                                            "email-index", receiverEmail)
    if record is None:
      return None
    
    # if there are multiple records, then match with any code

    for r in record:
    
      tempUserId = r['tempUserId'] if 'tempUserId' in r else None
      code = r['emailVerificationCode'] if 'emailVerificationCode' in r else None
      codeGenerationTime = r['verificationCodeGenerationTime'] if 'verificationCodeGenerationTime' in r else None
      isCodeVerified = r['verified'] if 'verified' in r else None

      if code is None or codeGenerationTime is None:
        continue

      dtCodeGenerated = datetime.strptime(codeGenerationTime, "%Y-%m-%d %H:%M:%S")
      minutesDiff = (datetime.now() - dtCodeGenerated).total_seconds() / 60

      # it's a match if it is generated within last 15 mins
      if code is not None and code == emailCode and \
          isCodeVerified != 'y' and minutesDiff <= 15:
        
        # update the table 
        DBManager.updateRecordInDynamoTable(DBTables.TempUser_Table_Name, "tempUserId",
                                                str(tempUserId), "email", receiverEmail, 
                                                {
                                                  "verified": "y",
                                                })

        return True
      else:
        continue
    
    return False

class PED_Module:
   
   @staticmethod
   def initiate():
      DBTables.GetTableName()
      Utility.initiate()

class DBTables:
   User_Table_Name = ""
   UserFiles_Table_Name = ""
   UserActivity_Table_Name = ""
   TempUser_Table_Name = ""

   @staticmethod
   def GetTableName():
      
      if Utility.ENVIRONMENT == "dev":
        DBTables.User_Table_Name = Utility.USER_TABLE_NAME   
        DBTables.UserFiles_Table_Name = Utility.USERFILES_TABLE_NAME   
        DBTables.UserActivity_Table_Name = Utility.USERACTIVITY_TABLE_NAME
        DBTables.TempUser_Table_Name = Utility.TEMPUSER_TABLE_NAME
      else:
        if Utility.ENVIRONMENT == "test":
          DBTables.User_Table_Name = Utility.TEST_USER_TABLE_NAME   
          DBTables.UserFiles_Table_Name = Utility.TEST_USERFILES_TABLE_NAME   
          DBTables.UserActivity_Table_Name = Utility.TEST_USERACTIVITY_TABLE_NAME
          DBTables.TempUser_Table_Name = Utility.TEST_TEMPUSER_TABLE_NAME

        else:
          DBTables.User_Table_Name = Utility.PROD_USER_TABLE_NAME   
          DBTables.UserFiles_Table_Name = Utility.PROD_USERFILES_TABLE_NAME   
          DBTables.UserActivity_Table_Name = Utility.PROD_USERACTIVITY_TABLE_NAME 
          DBTables.TempUser_Table_Name = Utility.PROD_TEMPUSER_TABLE_NAME
