from pptx import Presentation
import os
from pathlib import Path
import logging


def getTextFromPPTX(pptFilePath):

    # check for file existance
    isExist = os.path.exists(pptFilePath)
    if not isExist:    
        return None     

    #check for PPTX file
    fileExtension = Path(pptFilePath).suffix
    if fileExtension != ".pptx":
        return None

    text = ""

    try:
        # get every text section of every slide and 
        # get the text 
        prs = Presentation(pptFilePath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = text + "\n\n" + run.text
    
        return text
    
    except Exception as e:
        logging.error(e)
        return None