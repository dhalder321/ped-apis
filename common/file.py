from common.prompts import Prompt
from common.model import retryModelForOutputType
from common.globals import Utility
from pptx import Presentation
from docx import Document
import os, re, shutil
from pathlib import Path
import logging
import fitz
# from pdfminer.high_level import extract_pages
# from pdfminer.layout import LTTextContainer, LTChar


def getTextFromPPTX(pptFilePath):

    # check for file existance
    isExist = os.path.exists(pptFilePath)
    if not isExist:    
        return None     

    #check for PPTX file
    fileExtension = Path(pptFilePath).suffix
    if fileExtension != ".pptx":
        return None

    text = ""

    try:
        # get every text section of every slide and 
        # get the text 
        prs = Presentation(pptFilePath)
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text = text + "\n" + run.text
            #TODO: get the text from slide notes
    
        return text
    
    except Exception as e:
        logging.error(e)
        return None


def getTextFromDoc(docFilePath):

     # check for file existance
    isExist = os.path.exists(docFilePath)
    if not isExist:    
        return None     

    #check for docx file
    fileExtension = Path(docFilePath).suffix
    if fileExtension != ".docx":
        return None

    document = Document(docFilePath)
    full_text = []
    for paragraph in document.paragraphs:
        full_text.append(paragraph.text)
    return '\n'.join(full_text)

def getTextFromPDF(pdfFilePath):

    # check for file existance
    isExist = os.path.exists(pdfFilePath)
    if not isExist:    
        return None     

    #check for docx file
    fileExtension = Path(pdfFilePath).suffix
    if fileExtension != ".pdf":
        return None

    # Open the PDF file using fitz.open()
    doc = fitz.open(pdfFilePath)
    pageText = ""
    for page in range(doc.page_count):
        pageText += doc.load_page(page).get_text("text") + "\n"
    doc.close()

    return pageText

def getAudioScripts(pptFilePath):

    # check for file existance
    isExist = os.path.exists(pptFilePath)
    if not isExist:    
        return None     

    #check for PPTX file
    fileExtension = Path(pptFilePath).suffix
    if fileExtension != ".pptx":
        return None

    try:
        # check the notes section for each slide
        # if there is no notes, generate and set a note from the slide text
        # save the notes as audio script in the ./script folder

        scriptFileLocation = str(Path(Path(pptFilePath).parent, "script"))
        Path(scriptFileLocation).mkdir(parents=True, exist_ok=True)
        prs = Presentation(pptFilePath)
        slideCount = 1
        for slide in prs.slides:
            slideText = ''
            slideNote = ''
            if slide.notes_slide is not None and \
                slide.notes_slide.notes_text_frame is not None:

                slideNote = slide.notes_slide.notes_text_frame.text

            if slideNote is None or slideNote == '':
                # get the slide text
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            slideText = slideText + "\n" + run.text
            
                # generate and set a note
                prompt = Prompt.getPrompt(Utility.SLIDE_NOTE_FOR_VIDEO_PROMPT_TYPE) 
                if prompt is None:
                    raise Exception("Prompt could not be generated for slide note generation in checkOrModifyPPT4VideoGeneration method")
                
                prompt += "\n" + slideText

                slideNote = retryModelForOutputType(Utility.TRASFORMATION_USER_ROLE, prompt, 'text', maxRetry= 2)

                if slideNote is not None and slideNote != '':
                    slide.notes_slide.notes_text_frame.text = slideNote

            # save the node text in script folder
            scriptFileName = str(slideCount) + ".txt"
            scriptFilePath = Path(scriptFileLocation, scriptFileName)
            Path(scriptFilePath).touch(exist_ok=True)
            with scriptFilePath.open("w", encoding ="utf-8") as f:
                f.write(slideNote)

            slideCount += 1    

        # prs.Close()
    
        return str(scriptFilePath)
    
    except Exception as e:
        logging.error(e)
        return None


def getFolderSize(folder_path):
    total_size = 0
    for file in os.listdir(folder_path):
        file_path = os.path.join(folder_path, file)
        if os.path.isfile(file_path):
            total_size += os.path.getsize(file_path)
        elif os.path.isdir(file_path):
            total_size += getFolderSize(file_path)
    return total_size


def deleteDirWithFiles(folderPath):
    if folderPath is None or folderPath == '':
        return
    dir_path = Path(folderPath)
    if dir_path.exists():
        # Delete the directory and its contents
        for file in dir_path.rglob('*'):
            if file.is_file():
                file.unlink()
            elif file.is_dir():
                # file.rmdir()
                shutil.rmtree(file)
        dir_path.rmdir()
        print(f"Directory '{dir_path}' and its contents have been deleted.")
    else:
        print(f"The directory '{dir_path}' does not exist.")

# def get_text_coordinates(text_element):
#     """ Helper function to get the y-coordinate of a text element. """
#     return text_element.y0

# def filter_unwanted_text(text):
#     """Remove standalone dates and numbers which might represent page numbers or headers/footers."""
#     if re.match(r"^\d{1,2}/\d{1,2}/\d{2,4}$", text):  # Matches dates like 01/12/2021
#         return False
#     if text.isdigit() or (text.replace('.', '', 1).isdigit() and text.count('.') < 2):  # Check if it's a number or decimal
#         return False
#     return True

# def extract_text_with_headings(pdf_path, min_heading_font_size=12):
#     paragraphs = {}
#     current_heading = None
#     found_heading = False
#     text_blocks = []

#     for page_layout in extract_pages(pdf_path):
#         # Extract all text elements with their positions
#         for element in page_layout:
#             if isinstance(element, LTTextContainer):
#                 for text_line in element:
#                     text = text_line.get_text().strip()
#                     if text and filter_unwanted_text(text):
#                         bbox = text_line.bbox
#                         text_blocks.append((text, bbox[0], bbox[1]))  # x0, y0

#     # Sort blocks by y coordinate, then by x; this sorts top-to-bottom, left-to-right
#     text_blocks.sort(key=lambda b: (-b[2], b[1]))

#     # Concatenate text blocks into paragraphs while handling column layout
#     buffer = []
#     last_y = None
#     column_threshold = 50  # Adjust based on document layout specifics

#     for text, x, y in text_blocks:
#         if last_y is not None and abs(y - last_y) > column_threshold:
#             # Check if it's a new column or significant vertical jump
#             if buffer:
#                 full_text = ' '.join(buffer).replace('- ', '')  # Handle hyphenation
#                 # Determine if it's a heading
#                 chars = [obj for obj in page_layout if isinstance(obj, LTChar)]
#                 font_size = max((char.size for char in chars), default=0)
#                 if font_size >= min_heading_font_size and not found_heading:
#                     if current_heading:
#                         paragraphs[current_heading] = full_text
#                     current_heading = full_text
#                     found_heading = True
#                     buffer = []
#                 else:
#                     buffer.append(full_text)
#                 buffer = []
#         buffer.append(text)
#         last_y = y

#     # Save the last processed text as paragraph
#     if buffer:
#         full_text = ' '.join(buffer).replace('- ', '')
#         if current_heading and full_text:
#             paragraphs[current_heading] = full_text

#     return paragraphs if found_heading else ' '.join(buffer)  # Return as a single text if no headings found
